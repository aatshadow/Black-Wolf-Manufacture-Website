#!/usr/bin/env python3
"""
Generate BlackWolf Digitalization Bulgaria 2026 PPTX Presentation.
All text in Bulgarian. Dark premium corporate design.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Constants ──────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
BG_COLOR = RGBColor(0x0A, 0x0A, 0x1A)
ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB)
ACCENT_BLUE_LIGHT = RGBColor(0x3B, 0x82, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DIM_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_CARD = RGBColor(0x12, 0x12, 0x2A)
FONT = "Calibri"
LOGO_PATH = "/Users/alex/blackwolf-website/public/img/logo.png"
OUTPUT_PATH = "/Users/alex/blackwolf-website/presentations/blackwolf-digitalization-bg.pptx"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None,
              border_width=Pt(1), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.line.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txbox


def add_paragraph(text_frame, text, font_size=18, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(6)):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_decorative_bar(slide, left, top, width, height, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_title_with_bar(slide, title_text, subtitle_line=False):
    """Add a slide title with a blue accent bar on the left."""
    add_decorative_bar(slide, Inches(0.8), Inches(0.7), Inches(0.08), Inches(0.6), ACCENT_BLUE)
    add_text_box(slide, Inches(1.1), Inches(0.55), Inches(10), Inches(0.9),
                 title_text, font_size=38, color=WHITE, bold=True)
    if subtitle_line:
        add_decorative_bar(slide, Inches(1.1), Inches(1.45), Inches(2), Inches(0.03), ACCENT_BLUE)


def add_top_accent_line(slide):
    add_decorative_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.04), ACCENT_BLUE)


# ── Build Presentation ────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1 — Cover
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_COLOR)

# Decorative elements
add_decorative_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05), ACCENT_BLUE)
add_decorative_bar(slide, Inches(0), Inches(7.45), SLIDE_W, Inches(0.05), ACCENT_BLUE)

# Side accent bars
add_shape(slide, Inches(1.5), Inches(2.0), Inches(0.04), Inches(3.5),
          fill_color=ACCENT_BLUE, shape_type=MSO_SHAPE.RECTANGLE)
add_shape(slide, Inches(11.8), Inches(2.0), Inches(0.04), Inches(3.5),
          fill_color=ACCENT_BLUE, shape_type=MSO_SHAPE.RECTANGLE)

# Logo
if os.path.exists(LOGO_PATH):
    slide.shapes.add_picture(LOGO_PATH, Inches(5.4), Inches(1.0), Inches(2.5))

# Title
add_text_box(slide, Inches(2), Inches(3.0), Inches(9.333), Inches(0.8),
             "BLACKWOLF", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(2), Inches(3.9), Inches(9.333), Inches(0.7),
             "Цифрова Инфраструктура за Производство", font_size=26, color=ACCENT_BLUE_LIGHT,
             bold=False, alignment=PP_ALIGN.CENTER)

# Program line
add_text_box(slide, Inches(2), Inches(4.7), Inches(9.333), Inches(0.6),
             "Програма за Дигитализация — България 2026", font_size=18, color=DIM_GRAY,
             bold=False, alignment=PP_ALIGN.CENTER)

# Decorative bottom accent
add_decorative_bar(slide, Inches(5.5), Inches(5.6), Inches(2.3), Inches(0.03), ACCENT_BLUE)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2 — The Problem
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_COLOR)
add_top_accent_line(slide)
add_title_with_bar(slide, "Предизвикателството", subtitle_line=True)

problems = [
    ("73%", "на производителите работят с таблици"),
    ("€2.4M", "годишни загуби от неефективност"),
    ("0%", "свързаност — разпокъсани системи"),
    ("0", "видимост в реално време"),
    ("0", "киберсигурност"),
]

for i, (stat, desc) in enumerate(problems):
    y = Inches(2.0) + Inches(i * 1.0)
    # Card background
    card = add_shape(slide, Inches(1.1), y, Inches(10.5), Inches(0.85),
                     fill_color=DARK_CARD, border_color=RGBColor(0x1E, 0x1E, 0x3A))
    # Stat
    add_text_box(slide, Inches(1.4), y + Inches(0.15), Inches(1.5), Inches(0.55),
                 stat, font_size=28, color=ACCENT_BLUE, bold=True)
    # Description
    add_text_box(slide, Inches(3.0), y + Inches(0.15), Inches(8), Inches(0.55),
                 desc, font_size=20, color=LIGHT_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3 — Who We Are
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_COLOR)
add_top_accent_line(slide)
add_title_with_bar(slide, "Кои сме ние", subtitle_line=True)

who_items = [
    "BlackWolf — европейска компания за цифрова инфраструктура",
    "Специализирани в производствения сектор",
    "Екип с опит в NASA, Instagram, Uber Eats",
    "Офиси в Испания и България",
]

for i, item in enumerate(who_items):
    y = Inches(2.2) + Inches(i * 1.15)
    card = add_shape(slide, Inches(1.1), y, Inches(10.5), Inches(0.95),
                     fill_color=DARK_CARD, border_color=RGBColor(0x1E, 0x1E, 0x3A))
    # Blue dot / indicator
    add_decorative_bar(slide, Inches(1.3), y + Inches(0.35), Inches(0.25), Inches(0.25), ACCENT_BLUE)
    add_text_box(slide, Inches(1.8), y + Inches(0.2), Inches(9.5), Inches(0.6),
                 item, font_size=22, color=WHITE)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4 — What We Deliver
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_COLOR)
add_top_accent_line(slide)
add_title_with_bar(slide, "Какво предлагаме", subtitle_line=True)

services = [
    ("CRM система", "Управление на клиенти и продажби"),
    ("ERP система", "Управление на операции и инвентар"),
    ("BI платформа", "Бизнес анализи в реално време"),
    ("Уебсайтове и eCommerce", "Онлайн присъствие и продажби"),
    ("Киберсигурност (SOC)", "24/7 AI защита"),
]

card_w = Inches(3.2)
card_h = Inches(2.2)
gap = Inches(0.35)
# Row 1: 3 cards, Row 2: 2 cards centered
for i, (title, desc) in enumerate(services):
    if i < 3:
        x = Inches(1.1) + (card_w + gap) * i
        y = Inches(2.0)
    else:
        x = Inches(2.85) + (card_w + gap) * (i - 3)
        y = Inches(4.6)

    card = add_shape(slide, x, y, card_w, card_h,
                     fill_color=DARK_CARD, border_color=ACCENT_BLUE)
    # Blue top bar on card
    add_decorative_bar(slide, x + Inches(0.3), y + Inches(0.25), Inches(0.8), Inches(0.04), ACCENT_BLUE)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.5), card_w - Inches(0.6), Inches(0.5),
                 title, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.2), card_w - Inches(0.6), Inches(0.7),
                 desc, font_size=16, color=LIGHT_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5 — Process (5 weeks)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_COLOR)
add_top_accent_line(slide)
add_title_with_bar(slide, "Как работим — 5 седмици", subtitle_line=True)

weeks = [
    ("Седмица 1", "Анализ"),
    ("Седмица 2", "Архитектура"),
    ("Седмица 3", "Разработка"),
    ("Седмица 4", "Внедряване"),
    ("Седмица 5", "Оптимизация"),
]

step_w = Inches(2.0)
step_h = Inches(3.0)
total_w = step_w * 5 + Inches(0.3) * 4
start_x = (SLIDE_W - total_w) / 2

# Connecting line
line_y = Inches(3.5)
add_decorative_bar(slide, start_x + Inches(1.0), line_y, total_w - Inches(2.0), Inches(0.03), ACCENT_BLUE)

for i, (week, task) in enumerate(weeks):
    x = start_x + (step_w + Inches(0.3)) * i
    y = Inches(2.5)

    # Circle-like indicator
    circle = add_shape(slide, x + Inches(0.7), Inches(3.25), Inches(0.55), Inches(0.55),
                       fill_color=ACCENT_BLUE, border_color=ACCENT_BLUE,
                       shape_type=MSO_SHAPE.OVAL)
    # Number in circle
    add_text_box(slide, x + Inches(0.7), Inches(3.3), Inches(0.55), Inches(0.45),
                 str(i + 1), font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Week label above
    add_text_box(slide, x, Inches(2.3), step_w, Inches(0.5),
                 week, font_size=14, color=DIM_GRAY, alignment=PP_ALIGN.CENTER)

    # Task below
    card = add_shape(slide, x + Inches(0.15), Inches(4.2), step_w - Inches(0.3), Inches(1.3),
                     fill_color=DARK_CARD, border_color=RGBColor(0x1E, 0x1E, 0x3A))
    add_text_box(slide, x + Inches(0.15), Inches(4.5), step_w - Inches(0.3), Inches(0.7),
                 task, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6 — Results
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_COLOR)
add_top_accent_line(slide)
add_title_with_bar(slide, "Доказани Резултати", subtitle_line=True)

# Kingly card
card = add_shape(slide, Inches(1.1), Inches(2.0), Inches(5.0), Inches(2.8),
                 fill_color=DARK_CARD, border_color=ACCENT_BLUE)
add_text_box(slide, Inches(1.4), Inches(2.2), Inches(4.4), Inches(0.5),
             "Kingly", font_size=24, color=ACCENT_BLUE_LIGHT, bold=True)

kingly_stats = [
    "-73% време за поръчки",
    "0 загубени клиенти",
    "4 служители преразпределени",
]
for j, stat in enumerate(kingly_stats):
    add_text_box(slide, Inches(1.6), Inches(2.9 + j * 0.5), Inches(4.2), Inches(0.45),
                 "→  " + stat, font_size=17, color=LIGHT_GRAY)

# FBA Pro Academy card
card = add_shape(slide, Inches(6.6), Inches(2.0), Inches(5.0), Inches(2.8),
                 fill_color=DARK_CARD, border_color=ACCENT_BLUE)
add_text_box(slide, Inches(6.9), Inches(2.2), Inches(4.4), Inches(0.5),
             "FBA Pro Academy", font_size=24, color=ACCENT_BLUE_LIGHT, bold=True)

fba_stats = [
    "€500K+ приходи за 1 месец",
    "5 системи интегрирани",
]
for j, stat in enumerate(fba_stats):
    add_text_box(slide, Inches(7.1), Inches(2.9 + j * 0.5), Inches(4.2), Inches(0.45),
                 "→  " + stat, font_size=17, color=LIGHT_GRAY)

# Bottom stat banner
banner = add_shape(slide, Inches(1.1), Inches(5.3), Inches(10.5), Inches(1.2),
                   fill_color=ACCENT_BLUE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.8),
             "Над 80% подобрение на оперативната ефективност",
             font_size=26, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7 — Program Compliance
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_COLOR)
add_top_accent_line(slide)
add_title_with_bar(slide, "Съответствие с Програмата", subtitle_line=True)

checklist = [
    "ERP системи",
    "CRM системи",
    "Уебсайтове и eCommerce",
    "Бизнес интелигентност",
    "Киберсигурност",
    "AI и автоматизация",
]

# Two columns of 3
for i, item in enumerate(checklist):
    col = i // 3
    row = i % 3
    x = Inches(1.1) + col * Inches(5.5)
    y = Inches(2.2) + row * Inches(1.2)

    card = add_shape(slide, x, y, Inches(5.0), Inches(0.95),
                     fill_color=DARK_CARD, border_color=RGBColor(0x1E, 0x1E, 0x3A))
    # Checkmark indicator
    check_box = add_shape(slide, x + Inches(0.25), y + Inches(0.22), Inches(0.5), Inches(0.5),
                          fill_color=ACCENT_BLUE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.2), Inches(0.5), Inches(0.5),
                 "✓", font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(1.0), y + Inches(0.22), Inches(3.7), Inches(0.5),
                 item, font_size=20, color=WHITE)

# Bottom note
note_shape = add_shape(slide, Inches(3.0), Inches(6.0), Inches(7.0), Inches(0.8),
                       fill_color=DARK_CARD, border_color=ACCENT_BLUE)
add_text_box(slide, Inches(3.0), Inches(6.1), Inches(7.0), Inches(0.6),
             "Поемаме цялата документация",
             font_size=20, color=ACCENT_BLUE_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 8 — Why BlackWolf
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_COLOR)
add_top_accent_line(slide)
add_title_with_bar(slide, "Защо BlackWolf", subtitle_line=True)

why_blocks = [
    ("Всичко в едно", "CRM, ERP, BI, Уеб, SOC — един партньор"),
    ("5 седмици", "Бързо внедряване от анализ до оптимизация"),
    ("За производство", "Специализирани решения за индустрията"),
    ("AI технологии", "Изкуствен интелект вграден във всяко решение"),
]

card_w = Inches(5.0)
card_h = Inches(2.0)
gap_x = Inches(0.6)
gap_y = Inches(0.5)

for i, (title, desc) in enumerate(why_blocks):
    col = i % 2
    row = i // 2
    x = Inches(1.1) + col * (card_w + gap_x)
    y = Inches(2.2) + row * (card_h + gap_y)

    card = add_shape(slide, x, y, card_w, card_h,
                     fill_color=DARK_CARD, border_color=ACCENT_BLUE)
    add_decorative_bar(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.6), Inches(0.05), ACCENT_BLUE)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.5), card_w - Inches(0.6), Inches(0.5),
                 title, font_size=24, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.15), card_w - Inches(0.6), Inches(0.6),
                 desc, font_size=16, color=LIGHT_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 9 — Contact
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_COLOR)
add_decorative_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.05), ACCENT_BLUE)
add_decorative_bar(slide, Inches(0), Inches(7.45), SLIDE_W, Inches(0.05), ACCENT_BLUE)

# Side accent bars (matching cover)
add_shape(slide, Inches(1.5), Inches(1.5), Inches(0.04), Inches(4.5),
          fill_color=ACCENT_BLUE, shape_type=MSO_SHAPE.RECTANGLE)
add_shape(slide, Inches(11.8), Inches(1.5), Inches(0.04), Inches(4.5),
          fill_color=ACCENT_BLUE, shape_type=MSO_SHAPE.RECTANGLE)

# Logo
if os.path.exists(LOGO_PATH):
    slide.shapes.add_picture(LOGO_PATH, Inches(5.4), Inches(1.0), Inches(2.5))

add_text_box(slide, Inches(2), Inches(3.0), Inches(9.333), Inches(0.7),
             "Свържете се с нас", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_decorative_bar(slide, Inches(5.5), Inches(3.8), Inches(2.3), Inches(0.03), ACCENT_BLUE)

add_text_box(slide, Inches(2), Inches(4.1), Inches(9.333), Inches(0.5),
             "contact@blackwolfsec.io", font_size=22, color=ACCENT_BLUE_LIGHT,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(4.7), Inches(9.333), Inches(0.5),
             "blackwolfsec.io", font_size=22, color=ACCENT_BLUE_LIGHT,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.6), Inches(9.333), Inches(0.5),
             "Вашият партньор за дигитална трансформация", font_size=18, color=DIM_GRAY,
             alignment=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────
prs.save(OUTPUT_PATH)
print(f"Presentation saved to: {OUTPUT_PATH}")
print(f"Total slides: {len(prs.slides)}")
